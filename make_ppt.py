import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Please install python-pptx first: pip install python-pptx")
    exit(1)

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "EduGalaxy"
subtitle.text = "Where Learning Becomes Legend\nAn Immersive 3D Educational Simulator"

# Slide 2: The Problem
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Problem: Traditional Learning Lacks Engagement"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Attention spans are dropping faster than ever."
p = tf.add_paragraph()
p.text = "Standard syllabuses (CBSE/ICSE) are presented in dry, uninspiring formats."
p.level = 0
p = tf.add_paragraph()
p.text = "Students feel no sense of true progression or epic purpose in their daily homework."
p.level = 0

# Slide 3: The Solution
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Solution: EduGalaxy"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "A Cinematic Planetary Defense Simulator"
p = tf.add_paragraph()
p.text = "Students become 'Protagonist Heroes' defending the galaxy from an alien invasion."
p.level = 1
p = tf.add_paragraph()
p.text = "Completing the learning syllabus (Math, Science, English) powers the planetary shields."
p.level = 1
p = tf.add_paragraph()
p.text = "Includes interactive quizzes, puzzles, and 3D exploration labs."
p.level = 1

# Slide 4: Ultra-HD 3D Cinematic Experience
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Pushing the Web Beyond Limits"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Powered by WebGL and Three.js"
p = tf.add_paragraph()
p.text = "Massive interactive Star Wars-style galaxy with stunning 4K-style UI graphics."
p.level = 1
p = tf.add_paragraph()
p.text = "Features deep Z-axis scrolling, extreme mouse parallax, and massive glowing asteroids."
p.level = 1
p = tf.add_paragraph()
p.text = "Procedural generative audio engine (Tone.js) for evolving ambient music and laser SFX."
p.level = 1

# Slide 5: Mission Control & Leaderboards
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Tracking Legendary Progress"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Dynamic XP Tracking and Tier Leveling"
p = tf.add_paragraph()
p.text = "Session XP tracker bar with confetti celebrations for reaching daily goals."
p.level = 1
p = tf.add_paragraph()
p.text = "Achievement Rack with unlockable colorful SVG badges ('Star Cadet', 'Math Whiz')."
p.level = 1
p = tf.add_paragraph()
p.text = "Global competitive leaderboard filtered by Class (1-5)."
p.level = 1

# Slide 6: The Tech Stack
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Built for Speed and Scale"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Backend: Python & Flask for lightweight, lightning-fast API routing."
p = tf.add_paragraph()
p.text = "Frontend: Vanilla JS & CSS3 with Glassmorphism UI (Zero heavy UI frameworks for max performance)."
p.level = 0
p = tf.add_paragraph()
p.text = "3D & Audio: Three.js for heavy WebGL rendering and Tone.js for sound synthesis."
p.level = 0
p = tf.add_paragraph()
p.text = "Infrastructure: Fully offline-capable Progressive Web App (PWA) with a custom Service Worker."
p.level = 0

# Slide 7: Conclusion
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Let's Save the Galaxy"
subtitle.text = "Thank you!\nReady for a Live Demo?"

# Save to the documentation folder
output_path = r"c:\Users\AVROJIT\OneDrive\Desktop\EDUGALAXY\documentation\EduGalaxy_Pitch_Deck.pptx"
prs.save(output_path)
print(f"PPT successfully generated at: {output_path}")
